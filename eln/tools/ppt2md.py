#!/usr/bin/env python3
"""Convert .pptx to AI-readable markdown (one section per slide, extract text + speaker notes, export images).
Dependency: pip install python-pptx. Usage: python tools/ppt2md.py deck.pptx out_dir/"""
import sys, os
from pptx import Presentation
def convert(pptx, outdir):
    os.makedirs(outdir, exist_ok=True)
    prs=Presentation(pptx); md=[f"# {os.path.basename(pptx)}\n"]
    for i,slide in enumerate(prs.slides,1):
        md.append(f"\n## Slide {i}")
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                md.append(sh.text_frame.text.strip())
            if sh.shape_type==13:  # picture
                img=os.path.join(outdir,f"slide{i}_{sh.shape_id}.png")
                open(img,"wb").write(sh.image.blob); md.append(f"![](./{os.path.basename(img)})")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            md.append("**Speaker notes:** "+slide.notes_slide.notes_text_frame.text.strip())
    open(os.path.join(outdir,"slides.md"),"w").write("\n\n".join(md))
    print("->", os.path.join(outdir,"slides.md"))
if __name__=="__main__":
    convert(sys.argv[1], sys.argv[2])
